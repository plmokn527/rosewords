from pptx import Presentation
import pptx
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
import pandas as pd

数据 = pd.read_excel('test.xlsx')

标题 = []
内容 = []

for title in 数据['标题']:
    标题.append(title)
for content in 数据['内容']:
    内容.append(content)

# 打开幻灯片
prs = Presentation('test3.pptx')

# 获取要替换的幻灯片
slide = prs.slides[1]  # 假设要替换的是第一张幻灯片


def 修改文本框内容(第几条文本, shapes, 内容, 是否标题首页=False):
    index = 0
    if (第几条文本 == 1):
        index = 28;
    else:
        index = 第几条文本 - 2
    if (是否标题首页):
        index = 0
    textbox = shapes[index].text_frame
    textbox.clear()

    p = textbox.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()

    # 修改文本框内容
    run.text = f'{内容}'
    font = run.font
    font.name = '黑体'
    font.size = Pt(20)
    font.bold = False
    font.italic = None
    # font.color.theme_color = MSO_THEME_COLOR.ACCENT_6
    # 黑色
    font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def 修改标题框内容(第几条文本, shapes, 内容, 是否首页=False):
    index = 0
    if (第几条文本 == 1):
        index = 29;
    else:
        index = 第几条文本 - 2 + 14

    if (是否首页):
        index = 1
    textbox = shapes[index].text_frame
    textbox.clear()

    p = textbox.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()

    # 修改文本框内容
    run.text = f'{内容}'
    font = run.font
    font.name = '黑体'
    font.size = Pt(30)
    font.bold = True
    font.italic = None
    # font.color.theme_color = MSO_THEME_COLOR.ACCENT_6
    # 黑色
    font.color.rgb = RGBColor(0x00, 0x00, 0x00)


# 修改第二页
for shape in slide.shapes:
    # 修改图片
    if shape.shape_type == 6:
        print("group")
        # 获取组中的所有形状
        group_shapes = shape.shapes
        index = 30

        while (index < 36):
            new_pic_path = f'p{index - 29}.png'  # 新图片的路径
            print(new_pic_path)
            new_pptx_img = pptx.parts.image.Image.from_file(new_pic_path)
            # 获取要替换的图片
            old_pic = group_shapes[index]  # 假设要替换的是幻灯片中第一个形状的图片
            #
            index = index + 1
            # get part and rId from shape we need to change
            slide_part, rId = old_pic.part, old_pic._element.blip_rId
            image_part = slide_part.related_part(rId)

            # overwrite old blob info with new blob info
            image_part.blob = new_pptx_img._blob

        第几条内容文本 = 1
        while (第几条内容文本 < 16):
            修改文本框内容(第几条内容文本, group_shapes, 内容[第几条内容文本 - 1])
            第几条内容文本 = 第几条内容文本 + 1
        第几条标题文本 = 1
        while (第几条标题文本 < 16):
            修改标题框内容(第几条标题文本, group_shapes, 标题[第几条标题文本 - 1])
            第几条标题文本 = 第几条标题文本 + 1
# 修改第一页
slide = prs.slides[0]  # 假设要替换的是第一张幻灯片
group = 0
for shape in slide.shapes:
    # 修改图片
    if shape.shape_type == 6:
        group_shapes = shape.shapes
        new_pic_path = f'p{group + 1}.png'  # 新图片的路径
        new_pptx_img = pptx.parts.image.Image.from_file(new_pic_path)
        # 获取要替换的图片
        old_pic = group_shapes[2]  # 假设要替换的是幻灯片中第一个形状的图片
        #
        # get part and rId from shape we need to change
        slide_part, rId = old_pic.part, old_pic._element.blip_rId
        image_part = slide_part.related_part(rId)
        # overwrite old blob info with new blob info
        image_part.blob = new_pptx_img._blob

        # 修改标题
        第几条标题文本 = group + 1
        修改标题框内容(第几条标题文本, group_shapes, 标题[第几条标题文本 - 1], True)

        # 修改内容
        第几条内容文本 = group + 1
        修改文本框内容(第几条内容文本, group_shapes, 内容[第几条内容文本 - 1], True)

        group = group + 1

# 保存幻灯片
prs.save('test4.pptx')
